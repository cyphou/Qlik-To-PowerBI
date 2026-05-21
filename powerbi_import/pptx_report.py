"""PPTX report — generates PowerPoint executive summary.

Creates a migration summary presentation using python-pptx if
available, otherwise generates a Markdown fallback.
"""

from __future__ import annotations

import logging
import os
from typing import Any, Dict, List, Optional

logger = logging.getLogger('qlik_to_powerbi.pptx_report')

__all__ = ['PptxReport', 'generate_pptx_report']


class PptxReport:
    """Generates a PowerPoint migration summary."""

    def __init__(self, app_name: str = '', output_dir: str = '.'):
        self.app_name = app_name
        self.output_dir = output_dir

    def _build_slides(self, report_data: Dict[str, Any],
                      qa_data: Optional[Dict] = None) -> List[Dict[str, Any]]:
        """Build slide data structures."""
        slides: List[Dict[str, Any]] = []
        summary = report_data.get('summary', {})
        total = summary.get('total', 0) or 1
        exact = summary.get('exact', 0)
        pct = round(exact * 100 / total)

        # Title slide
        slides.append({
            'layout': 'title',
            'title': f'Qlik → Power BI Migration',
            'subtitle': f'{self.app_name} — Executive Summary',
        })

        # Fidelity overview
        slides.append({
            'layout': 'content',
            'title': 'Migration Fidelity',
            'content': [
                f'Overall fidelity: {pct}%',
                f'Total items: {total}',
                f'Exact matches: {exact}',
                f'Approximate: {summary.get("approximate", 0)}',
                f'Unsupported: {summary.get("unsupported", 0)}',
            ],
        })

        # Breakdown by type
        by_type = report_data.get('by_type', {})
        if by_type:
            content = []
            for item_type, counts in by_type.items():
                if isinstance(counts, dict):
                    content.append(f'{item_type}: {counts.get("total", 0)} items')
                else:
                    content.append(f'{item_type}: {counts}')
            slides.append({
                'layout': 'content',
                'title': 'Breakdown by Type',
                'content': content[:10],
            })

        # QA results
        if qa_data:
            qa_summary = qa_data.get('summary', {})
            slides.append({
                'layout': 'content',
                'title': 'QA Pipeline Results',
                'content': [
                    f'Total checks: {qa_summary.get("total_checks", 0)}',
                    f'Passed: {qa_summary.get("passed", 0)}',
                    f'Auto-fixed: {qa_summary.get("autofixed", 0)}',
                    f'Remaining issues: {qa_summary.get("remaining", 0)}',
                ],
            })

        # Recommendations
        items = report_data.get('items', [])
        unsupported = [i for i in items if i.get('status') == 'unsupported']
        if unsupported:
            slides.append({
                'layout': 'content',
                'title': 'Items Requiring Manual Review',
                'content': [
                    f'{i.get("name", "?")} ({i.get("type", "?")})' 
                    for i in unsupported[:8]
                ],
            })

        # Next steps
        slides.append({
            'layout': 'content',
            'title': 'Next Steps',
            'content': [
                'Review unsupported items for manual conversion',
                'Validate data accuracy in parallel run',
                'Configure data refresh schedules',
                'Set up row-level security',
                'Train end users on Power BI',
                'Plan cutover timeline',
            ],
        })

        return slides

    def generate_pptx(self, report_data: Dict[str, Any],
                      qa_data: Optional[Dict] = None) -> str:
        """Generate PPTX if python-pptx available, else Markdown fallback."""
        slides = self._build_slides(report_data, qa_data)
        os.makedirs(self.output_dir, exist_ok=True)

        try:
            from pptx import Presentation  # type: ignore[import-untyped]
            from pptx.util import Inches, Pt  # type: ignore[import-untyped]

            prs = Presentation()
            for slide_data in slides:
                if slide_data['layout'] == 'title':
                    layout = prs.slide_layouts[0]
                    slide = prs.slides.add_slide(layout)
                    slide.shapes.title.text = slide_data['title']
                    if slide.placeholders[1]:
                        slide.placeholders[1].text = slide_data.get('subtitle', '')
                else:
                    layout = prs.slide_layouts[1]
                    slide = prs.slides.add_slide(layout)
                    slide.shapes.title.text = slide_data['title']
                    body = slide.placeholders[1]
                    tf = body.text_frame
                    tf.clear()
                    for i, item in enumerate(slide_data.get('content', [])):
                        if i == 0:
                            tf.text = item
                        else:
                            p = tf.add_paragraph()
                            p.text = item

            filename = f'{self.app_name or "migration"}_summary.pptx'
            output_path = os.path.join(self.output_dir, filename)
            prs.save(output_path)
            logger.info("PPTX report saved to %s", output_path)
            return output_path

        except Exception:
            logger.warning("python-pptx not available or broken; generating Markdown")
            return self._generate_markdown(slides)

    def _generate_markdown(self, slides: List[Dict[str, Any]]) -> str:
        """Generate Markdown fallback."""
        lines = []
        for slide in slides:
            if slide['layout'] == 'title':
                lines.extend([
                    f'# {slide["title"]}',
                    f'*{slide.get("subtitle", "")}*',
                    '',
                ])
            else:
                lines.append(f'## {slide["title"]}')
                for item in slide.get('content', []):
                    lines.append(f'- {item}')
                lines.append('')

        filename = f'{self.app_name or "migration"}_summary.md'
        output_path = os.path.join(self.output_dir, filename)
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write('\n'.join(lines))
        logger.info("Markdown report saved to %s", output_path)
        return output_path


def generate_pptx_report(app_name: str, output_dir: str,
                         report_data: Dict[str, Any],
                         qa_data: Optional[Dict] = None) -> str:
    """Convenience function to generate a PPTX report."""
    report = PptxReport(app_name=app_name, output_dir=output_dir)
    return report.generate_pptx(report_data, qa_data)
