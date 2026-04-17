#!/usr/bin/env python3
"""Generate a PPTX presentation for the Qlik to Power BI Migration Tool."""

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Image paths ──
_SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
_PROJECT_ROOT = os.path.dirname(_SCRIPT_DIR)
IMG_MIGRATION_RESULTS = os.path.join(_PROJECT_ROOT, "docs", "images", "migration_results.png")
IMG_SHARE_ASSESSMENT = os.path.join(_PROJECT_ROOT, "docs", "images", "share_assessment.png")

# ── Brand colors ──
QLIK_GREEN = RGBColor(0x00, 0x98, 0x48)
PBI_YELLOW = RGBColor(0xF2, 0xC8, 0x11)
AZURE_BLUE = RGBColor(0x00, 0x78, 0xD4)
DARK_BG = RGBColor(0x1E, 0x1E, 0x2E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY = RGBColor(0x88, 0x88, 0x88)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
ACCENT_PURPLE = RGBColor(0x63, 0x66, 0xF1)


def _add_background(slide, color):
    """Set solid background color on a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_shape_bg(slide, color, left=0, top=0, width=None, height=None):
    """Add a colored rectangle behind content."""
    w = width or Inches(13.333)
    h = height or Inches(7.5)
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_accent_bar(slide, color, left, top, width, height):
    """Add a thin accent bar."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_textbox(slide, left, top, width, height, text, font_size=18,
                 color=DARK_TEXT, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name='Segoe UI'):
    """Add a text box with styling."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_bullet_list(slide, left, top, width, height, items, font_size=16,
                     color=DARK_TEXT, spacing=Pt(6)):
    """Add a bulleted list."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = 'Segoe UI'
        p.space_after = spacing
        p.level = 0
    return txBox


def _add_card(slide, left, top, width, height, title, body, accent_color=AZURE_BLUE):
    """Add a card-style box with accent bar."""
    # Card background
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    card.line.width = Pt(1)
    # Accent bar
    _add_accent_bar(slide, accent_color, left + Inches(0.15), top + Inches(0.15),
                    Inches(0.08), height - Inches(0.3))
    # Title
    _add_textbox(slide, left + Inches(0.4), top + Inches(0.1), width - Inches(0.5),
                 Inches(0.4), title, font_size=14, bold=True, color=accent_color)
    # Body
    _add_textbox(slide, left + Inches(0.4), top + Inches(0.45), width - Inches(0.5),
                 height - Inches(0.55), body, font_size=12, color=DARK_TEXT)


def _add_stat_card(slide, left, top, width, height, number, label, color=AZURE_BLUE):
    """Add a big-number stat card."""
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = color
    card.line.fill.background()
    _add_textbox(slide, left, top + Inches(0.15), width, Inches(0.6),
                 number, font_size=32, bold=True, color=WHITE, alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, left, top + Inches(0.7), width, Inches(0.4),
                 label, font_size=12, color=WHITE, alignment=PP_ALIGN.CENTER)


def _add_code_box(slide, left, top, width, height, code_text):
    """Add a code-style text box."""
    box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)
    box.line.fill.background()
    txBox = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15),
                                     width - Inches(0.4), height - Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(code_text.strip().split('\n')):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(11)
        p.font.color.rgb = RGBColor(0xA6, 0xE2, 0x2E)
        p.font.name = 'Cascadia Code'
        p.space_after = Pt(2)


def _section_header(slide, title, subtitle=''):
    """Standard section header layout."""
    _add_shape_bg(slide, LIGHT_GRAY)
    _add_accent_bar(slide, AZURE_BLUE, Inches(0.8), Inches(1.5), Inches(0.1), Inches(1.0))
    _add_textbox(slide, Inches(1.1), Inches(1.5), Inches(10), Inches(0.7),
                 title, font_size=36, bold=True, color=DARK_TEXT)
    if subtitle:
        _add_textbox(slide, Inches(1.1), Inches(2.2), Inches(10), Inches(0.5),
                     subtitle, font_size=18, color=MID_GRAY)


def build_presentation():
    """Build and save the PPTX."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]  # blank layout

    # ══════════════════════════════════════════════════════════════
    # SLIDE 1 — Title
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _add_shape_bg(slide, DARK_BG)
    # Gradient accent stripe
    _add_accent_bar(slide, QLIK_GREEN, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06))
    _add_accent_bar(slide, PBI_YELLOW, Inches(0), Inches(3.28), Inches(13.333), Inches(0.06))
    # Title
    _add_textbox(slide, Inches(1), Inches(1.8), Inches(11), Inches(1.0),
                 'Qlik to Power BI Migration Tool', font_size=44, bold=True,
                 color=WHITE, alignment=PP_ALIGN.LEFT)
    _add_textbox(slide, Inches(1), Inches(2.6), Inches(11), Inches(0.5),
                 'Automated Qlik Sense migration — Input, Output & Usage Guide',
                 font_size=20, color=PBI_YELLOW, alignment=PP_ALIGN.LEFT)
    # Bottom info
    _add_textbox(slide, Inches(1), Inches(5.5), Inches(6), Inches(0.4),
                 'Python 3.12+ · Zero external dependencies · 2,000 tests',
                 font_size=14, color=MID_GRAY)
    _add_textbox(slide, Inches(1), Inches(5.9), Inches(6), Inches(0.4),
                 'v9.1.0 — Preceptorship Multi-Agent Model', font_size=14, color=MID_GRAY)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 2 — What It Does
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _add_shape_bg(slide, LIGHT_GRAY)
    _section_header(slide, 'What Does This Tool Do?')
    _add_textbox(slide, Inches(1.1), Inches(3.0), Inches(10.5), Inches(0.8),
                 'Converts Qlik Sense applications (.qvf / JSON exports) into '
                 'complete Power BI projects (.pbip) with PBIR v4.0 reports and TMDL semantic models — '
                 'ready to open in Power BI Desktop or deploy to Fabric.',
                 font_size=18, color=DARK_TEXT)
    # 3 value prop cards
    _add_card(slide, Inches(1.1), Inches(4.2), Inches(3.4), Inches(2.2),
              '🔍  Automated Extraction',
              '11 JSON intermediates extracted from Qlik QVF/JSON: '
              'app metadata, datasources, dimensions, measures, '
              'visualizations, sheets, variables, load script, '
              'associations, bookmarks, master items.',
              QLIK_GREEN)
    _add_card(slide, Inches(4.9), Inches(4.2), Inches(3.4), Inches(2.2),
              '⚙️  Smart Generation',
              '175+ DAX formula conversions, 75+ visual type mappings, 40+ Power Query M transforms, '
              'auto-Calendar table, RLS roles, hierarchies, themes, set analysis → CALCULATE.',
              AZURE_BLUE)
    _add_card(slide, Inches(8.7), Inches(4.2), Inches(3.4), Inches(2.2),
              '🚀  Deploy Anywhere',
              'Deploy to Power BI Service, Microsoft Fabric, or Docker REST API. '
              'Bundle deployment, multi-tenant, blue/green rolling deploy with auto-rollback.',
              GREEN)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 3 — Architecture / Pipeline
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _add_shape_bg(slide, WHITE)
    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
                 '2-Step Pipeline Architecture', font_size=32, bold=True, color=DARK_TEXT)
    _add_textbox(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.5),
                 '.qvf → [Extraction] → 11 JSON files → [Generation] → .pbip (PBIR + TMDL)',
                 font_size=16, color=MID_GRAY, font_name='Cascadia Code')

    # Step 1 box
    _add_card(slide, Inches(0.6), Inches(1.8), Inches(5.5), Inches(4.8),
              'Step 1 — EXTRACT  (qlik_export/)',
              'Parses Qlik QVF/JSON into 11 intermediate JSON files.\n\n'
              '• extraction_orchestrator.py — main orchestrator\n'
              '• qvf_extractor.py — .qvf ZIP reader\n'
              '• dax_converter.py — 175+ formula conversions\n'
              '• m_query_builder.py — 25 connectors + 40+ transforms\n'
              '• qlik_script_converter.py — load script → M\n'
              '• format_adapter.py — bridge to generation layer\n'
              '• datasource_extractor.py — type/formula adapters',
              QLIK_GREEN)

    # Arrow
    _add_textbox(slide, Inches(6.2), Inches(3.5), Inches(0.8), Inches(0.6),
                 '→', font_size=48, bold=True, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # Step 2 box
    _add_card(slide, Inches(7.0), Inches(1.8), Inches(5.5), Inches(4.8),
              'Step 2 — GENERATE  (powerbi_import/)',
              'Converts JSON into a complete .pbip project.\n\n'
              '• tmdl_generator.py — TMDL semantic model\n'
              '• pbip_generator.py — .pbip project structure\n'
              '• visual_generator.py — 75+ visual type mappings\n'
              '• api_server.py — REST API server\n'
              '• schema_drift.py — drift detection\n'
              '• assessment.py — readiness scoring\n'
              '• shared_model.py — multi-app merge\n'
              '• deploy/ — Fabric & PBI Service deployment',
              AZURE_BLUE)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 4 — Input Formats
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _add_shape_bg(slide, LIGHT_GRAY)
    _section_header(slide, 'Supported Inputs', 'What you can feed into the tool')

    # Input cards
    _add_card(slide, Inches(0.6), Inches(3.2), Inches(3.8), Inches(3.5),
              '📄  .qvf',
              'Qlik Sense Application\n(ZIP archive)\n\nContains app metadata, data model, '
              'sheets, visualizations, variables, load script, bookmarks.',
              QLIK_GREEN)
    _add_card(slide, Inches(4.8), Inches(3.2), Inches(3.8), Inches(3.5),
              '📦  JSON Export',
              'Engine API Export\n(JSON files)\n\nMetadata exports from Qlik Sense Engine '
              'API — app properties, dimensions, measures, sheets.',
              QLIK_GREEN)
    _add_card(slide, Inches(9.0), Inches(3.2), Inches(3.8), Inches(3.5),
              '🔄  Two-Step Workflow',
              'Extract → Generate\n\nExtract once, generate multiple times.\n'
              'Use --skip-extraction to reuse\npreviously extracted JSON files.\n'
              'Ideal for iterative tuning.',
              AZURE_BLUE)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 5 — 17 Extracted Objects
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _add_shape_bg(slide, WHITE)
    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
                 '11 Intermediate JSON Files', font_size=32, bold=True, color=DARK_TEXT)
    _add_textbox(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
                 'Each extracted to a structured JSON file for the generation step',
                 font_size=16, color=MID_GRAY)

    objects = [
        ('app_metadata', 'App name, description, author, dates'),
        ('datasources', 'Connection strings, tables, columns, types'),
        ('dimensions', 'Master dimensions (fields, labels, groupings)'),
        ('measures', 'Master measures (expressions, labels, formats)'),
        ('visualizations', 'Chart types, dimension/measure bindings'),
        ('sheets', 'Sheet layouts, cell positions'),
        ('variables', 'Variables (name, definition, comment)'),
        ('loadscript', 'Full Qlik load script'),
        ('associations', 'Table associations / relationships'),
        ('bookmarks', 'Bookmarks and selections'),
        ('master_items', 'Master items (combined dim/measure refs)'),
    ]
    col1 = objects[:6]
    col2 = objects[6:]
    for i, (name, desc) in enumerate(col1):
        y = Inches(1.5) + Inches(i * 0.6)
        _add_textbox(slide, Inches(0.8), y, Inches(2.2), Inches(0.5),
                     name + '.json', font_size=13, bold=True, color=AZURE_BLUE,
                     font_name='Cascadia Code')
        _add_textbox(slide, Inches(3.0), y, Inches(3.5), Inches(0.5),
                     desc, font_size=13, color=DARK_TEXT)
    for i, (name, desc) in enumerate(col2):
        y = Inches(1.5) + Inches(i * 0.6)
        _add_textbox(slide, Inches(7.0), y, Inches(2.2), Inches(0.5),
                     name + '.json', font_size=13, bold=True, color=AZURE_BLUE,
                     font_name='Cascadia Code')
        _add_textbox(slide, Inches(9.2), y, Inches(3.5), Inches(0.5),
                     desc, font_size=13, color=DARK_TEXT)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 6 — Output Formats
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _add_shape_bg(slide, LIGHT_GRAY)
    _section_header(slide, 'Output Formats', 'Two output modes for different targets')

    # PBIP output
    _add_card(slide, Inches(0.6), Inches(3.2), Inches(5.8), Inches(3.8),
              '📊  Default: .pbip Project (PBIR v4.0 + TMDL)',
              '  ProjectName/\n'
              '  ├── ProjectName.pbip\n'
              '  ├── ProjectName.Report/\n'
              '  │   ├── report.json\n'
              '  │   └── definition/pages/.../visual.json\n'
              '  └── ProjectName.SemanticModel/\n'
              '      └── definition/\n'
              '          ├── model.tmdl, relationships.tmdl\n'
              '          ├── roles.tmdl, perspectives.tmdl\n'
              '          └── tables/*.tmdl\n\n'
              'Opens directly in Power BI Desktop (Dec 2025+).',
              AZURE_BLUE)

    # Fabric output
    _add_card(slide, Inches(6.8), Inches(3.2), Inches(5.8), Inches(3.8),
              '🏭  Fabric: --output-format fabric',
              '  5 Fabric artifacts generated:\n\n'
              '  1. Lakehouse — Delta table schemas + DDL\n'
              '  2. Dataflow Gen2 — Power Query M ingestion\n'
              '  3. PySpark Notebook — ETL transformations\n'
              '  4. DirectLake Semantic Model — TMDL\n'
              '  5. Data Pipeline — 3-stage orchestration\n\n'
              'Deploys to Microsoft Fabric workspaces.',
              RGBColor(0x00, 0x78, 0xD4))

    # ══════════════════════════════════════════════════════════════
    # SLIDE 7 — Basic CLI Usage
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _add_shape_bg(slide, WHITE)
    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
                 'Basic Usage — CLI Commands', font_size=32, bold=True, color=DARK_TEXT)

    _add_textbox(slide, Inches(0.8), Inches(1.2), Inches(5), Inches(0.4),
                 'Single QVF migration', font_size=16, bold=True, color=AZURE_BLUE)
    _add_code_box(slide, Inches(0.8), Inches(1.6), Inches(11.5), Inches(0.6),
                  'python migrate.py your_app.qvf')

    _add_textbox(slide, Inches(0.8), Inches(2.4), Inches(5), Inches(0.4),
                 'From JSON export', font_size=16, bold=True, color=AZURE_BLUE)
    _add_code_box(slide, Inches(0.8), Inches(2.8), Inches(11.5), Inches(0.6),
                  'python migrate.py export.json --output-dir output/my_project --verbose')

    _add_textbox(slide, Inches(0.8), Inches(3.6), Inches(5), Inches(0.4),
                 'Two-step (reuse extracted JSON)', font_size=16, bold=True, color=AZURE_BLUE)
    _add_code_box(slide, Inches(0.8), Inches(4.0), Inches(11.5), Inches(0.6),
                  'python migrate.py app.qvf --skip-extraction --output-dir output/step1')

    _add_textbox(slide, Inches(0.8), Inches(4.8), Inches(5), Inches(0.4),
                 'Interactive wizard', font_size=16, bold=True, color=AZURE_BLUE)
    _add_code_box(slide, Inches(0.8), Inches(5.2), Inches(11.5), Inches(0.6),
                  'python migrate.py app.qvf --wizard')

    _add_textbox(slide, Inches(0.8), Inches(6.0), Inches(5), Inches(0.4),
                 'Fabric-native output', font_size=16, bold=True, color=AZURE_BLUE)
    _add_code_box(slide, Inches(0.8), Inches(6.4), Inches(11.5), Inches(0.6),
                  'python migrate.py app.qvf --output-format fabric')

    # ══════════════════════════════════════════════════════════════
    # SLIDE 8 — Advanced CLI
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _add_shape_bg(slide, WHITE)
    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
                 'Advanced CLI Flags', font_size=32, bold=True, color=DARK_TEXT)

    flags = [
        ('--assess', 'Run pre-migration readiness assessment'),
        ('--dry-run', 'Preview migration without writing files'),
        ('--json', 'Machine-readable JSON output for CI/CD'),
        ('--merge app1.json app2.json', 'Merge apps into shared semantic model'),
        ('--deploy WORKSPACE_ID', 'Deploy to Power BI Service'),
        ('--deploy-refresh', 'Trigger dataset refresh after deployment'),
        ('--assess-server DIR', 'Portfolio-level server assessment'),
        ('--output-format fabric', 'Generate Fabric-native artifacts'),
        ('--plugins module.Class', 'Load custom migration plugins'),
        ('--optimize-dax', 'Run DAX optimizer (IF→SWITCH, COALESCE)'),
        ('--verbose', 'Detailed logging output'),
        ('--skip-extraction', 'Reuse previously extracted JSON'),
    ]
    for i, (flag, desc) in enumerate(flags):
        y = Inches(1.2) + Inches(i * 0.48)
        _add_textbox(slide, Inches(0.8), y, Inches(4.5), Inches(0.4),
                     flag, font_size=13, bold=True, color=AZURE_BLUE,
                     font_name='Cascadia Code')
        _add_textbox(slide, Inches(5.5), y, Inches(7), Inches(0.4),
                     desc, font_size=14, color=DARK_TEXT)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 9 — Qlik Server Integration
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _add_shape_bg(slide, LIGHT_GRAY)
    _section_header(slide, 'REST API & Docker',
                    'Headless / programmatic migration via HTTP')

    _add_textbox(slide, Inches(1.1), Inches(3.2), Inches(5), Inches(0.4),
                 'Start the REST API server', font_size=16, bold=True, color=AZURE_BLUE)
    _add_code_box(slide, Inches(1.1), Inches(3.7), Inches(11), Inches(0.9),
                  '# Direct\n'
                  'python -m powerbi_import.api_server\n\n'
                  '# Docker\n'
                  'docker build -t qlik-to-pbi . && docker run -p 8000:8000 qlik-to-pbi')

    _add_textbox(slide, Inches(1.1), Inches(4.9), Inches(5), Inches(0.4),
                 'Upload a QVF for migration', font_size=16, bold=True, color=AZURE_BLUE)
    _add_code_box(slide, Inches(1.1), Inches(5.4), Inches(11), Inches(0.7),
                  'curl -X POST -F "file=@app.qvf" http://localhost:8000/migrate\n'
                  'curl http://localhost:8000/status/{job_id}\n'
                  'curl -o result.zip http://localhost:8000/download/{job_id}')

    # ══════════════════════════════════════════════════════════════
    # SLIDE 10 — REST API
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _add_shape_bg(slide, WHITE)
    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
                 'REST API Server', font_size=32, bold=True, color=DARK_TEXT)
    _add_textbox(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
                 'Headless / programmatic migration via HTTP  (Sprint 110)',
                 font_size=16, color=MID_GRAY)

    _add_textbox(slide, Inches(0.8), Inches(1.6), Inches(5), Inches(0.4),
                 'Start the server', font_size=16, bold=True, color=AZURE_BLUE)
    _add_code_box(slide, Inches(0.8), Inches(2.0), Inches(5.5), Inches(1.0),
                  '# Direct\n'
                  'python -m powerbi_import.api_server\n\n'
                  '# Docker\n'
                  'docker build -t Qlik-to-pbi .\n'
                  'docker run -p 8000:8000 Qlik-to-pbi')

    # Endpoints table
    endpoints = [
        ('POST', '/migrate', 'Upload .qvf/.qvf for migration'),
        ('GET', '/status/{id}', 'Check job status'),
        ('GET', '/download/{id}', 'Download .pbip as ZIP'),
        ('GET', '/health', 'Health check'),
        ('GET', '/jobs', 'List all jobs'),
    ]
    _add_textbox(slide, Inches(7.0), Inches(1.6), Inches(5), Inches(0.4),
                 'Endpoints', font_size=16, bold=True, color=AZURE_BLUE)
    for i, (method, path, desc) in enumerate(endpoints):
        y = Inches(2.1) + Inches(i * 0.45)
        _add_textbox(slide, Inches(7.0), y, Inches(0.8), Inches(0.4),
                     method, font_size=13, bold=True, color=GREEN, font_name='Cascadia Code')
        _add_textbox(slide, Inches(7.8), y, Inches(2.0), Inches(0.4),
                     path, font_size=13, color=DARK_TEXT, font_name='Cascadia Code')
        _add_textbox(slide, Inches(9.9), y, Inches(3.0), Inches(0.4),
                     desc, font_size=13, color=MID_GRAY)

    # Usage example
    _add_textbox(slide, Inches(0.8), Inches(3.3), Inches(5), Inches(0.4),
                 'Usage example', font_size=16, bold=True, color=AZURE_BLUE)
    _add_code_box(slide, Inches(0.8), Inches(3.8), Inches(11.5), Inches(2.8),
                  '# 1. Upload a workbook\n'
                  'curl -X POST -F "file=@workbook.qvf" http://localhost:8000/migrate\n'
                  '# → {"job_id": "abc123", "status": "queued"}\n\n'
                  '# 2. Check status\n'
                  'curl http://localhost:8000/status/abc123\n'
                  '# → {"job_id": "abc123", "status": "completed"}\n\n'
                  '# 3. Download result\n'
                  'curl -o output.zip http://localhost:8000/download/abc123')

    # ══════════════════════════════════════════════════════════════
    # SLIDE 11 — Assessment & Planning
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _add_shape_bg(slide, LIGHT_GRAY)
    _section_header(slide, 'Pre-Migration Assessment',
                    'Analyze complexity before migrating')

    _add_card(slide, Inches(0.6), Inches(3.2), Inches(3.8), Inches(3.5),
              '📋  Workbook Assessment',
              '9 assessment categories:\n\n'
              '• Datasource complexity\n'
              '• Calculation coverage\n'
              '• Visual type mapping\n'
              '• Filter complexity\n'
              '• Data model structure\n'
              '• Interactivity (actions)\n'
              '• Extract / Hyper data\n'
              '• Migration scope\n'
              '• Connection string audit',
              AZURE_BLUE)

    _add_card(slide, Inches(4.8), Inches(3.2), Inches(3.8), Inches(3.5),
              '🏢  Server Portfolio',
              'Portfolio-level planning:\n\n'
              '• Per-workbook GREEN/YELLOW/RED\n'
              '• 8-axis complexity scoring\n'
              '• Effort estimation\n'
              '• Migration wave planning\n'
              '• Connector census\n'
              '• HTML executive dashboard',
              ACCENT_PURPLE)

    _add_card(slide, Inches(9.0), Inches(3.2), Inches(3.8), Inches(3.5),
              '🎯  Strategy Advisor',
              'Recommends Import / DirectQuery / Composite\nbased on 7 signals:\n\n'
              '• Data volume\n'
              '• Refresh frequency\n'
              '• Concurrent users\n'
              '• Cross-source joins\n'
              '• Real-time needs\n'
              '• Data sensitivity\n'
              '• Licensing constraints',
              GREEN)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 12 — Shared Semantic Model
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _add_shape_bg(slide, WHITE)
    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
                 'Shared Semantic Model', font_size=32, bold=True, color=DARK_TEXT)
    _add_textbox(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
                 'Merge multiple Qlik workbooks into one unified Power BI data model',
                 font_size=16, color=MID_GRAY)

    # Flow diagram as text
    _add_code_box(slide, Inches(0.8), Inches(1.7), Inches(11.5), Inches(2.0),
                  '  Workbook A ──→ Extract ──→ 17 JSON (A)  ──┐\n'
                  '  Workbook B ──→ Extract ──→ 17 JSON (B)  ──┤── MERGE ──→ Shared SemanticModel\n'
                  '  Workbook C ──→ Extract ──→ 17 JSON (C)  ──┘       ├──→ Thin Report A\n'
                  '                                                     ├──→ Thin Report B\n'
                  '                                                     └──→ Thin Report C')

    _add_textbox(slide, Inches(0.8), Inches(4.0), Inches(5), Inches(0.4),
                 'Commands', font_size=16, bold=True, color=AZURE_BLUE)
    _add_code_box(slide, Inches(0.8), Inches(4.4), Inches(11.5), Inches(2.2),
                  '# Merge apps into shared model\n'
                  'python migrate.py --merge app1.json app2.json app3.json\n\n'
                  '# With Fabric-native output\n'
                  'python migrate.py --merge app1.json app2.json --output-format fabric\n\n'
                  '# Deploy to Power BI Service\n'
                  'python migrate.py --merge app1.json app2.json --deploy WORKSPACE_ID')

    # Share assessment screenshot
    if os.path.isfile(IMG_SHARE_ASSESSMENT):
        slide.shapes.add_picture(IMG_SHARE_ASSESSMENT,
                                 Inches(0.8), Inches(6.8), Inches(11.5))

    # ══════════════════════════════════════════════════════════════
    # SLIDE 13 — Schema Drift & Incremental
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _add_shape_bg(slide, LIGHT_GRAY)
    _section_header(slide, 'Schema Drift Detection',
                    'Track changes between migration runs  (Sprint 111)')

    _add_textbox(slide, Inches(1.1), Inches(3.0), Inches(10.5), Inches(0.5),
                 'Compare the current Qlik source against a previously saved extraction snapshot '
                 'to detect added, removed, or modified objects.',
                 font_size=16, color=DARK_TEXT)

    _add_code_box(slide, Inches(1.1), Inches(3.8), Inches(10.5), Inches(1.3),
                  '# First run: creates baseline snapshot\n'
                  'python migrate.py workbook.qvf --check-drift ./snapshots/sales\n\n'
                  '# Later: detect drift vs baseline\n'
                  'python migrate.py workbook.qvf --check-drift ./snapshots/sales')

    _add_textbox(slide, Inches(1.1), Inches(5.4), Inches(10.5), Inches(0.4),
                 'Detects changes in 7 categories:', font_size=14, bold=True, color=AZURE_BLUE)
    _add_bullet_list(slide, Inches(1.1), Inches(5.8), Inches(10.5), Inches(1.5),
                     ['Tables (added / removed)',
                      'Columns (added / removed / type changed)',
                      'Calculations (added / formula modified)',
                      'Worksheets (added / fields changed)',
                      'Relationships (added / removed)',
                      'Parameters (added / value changed)',
                      'Filters (added)'],
                     font_size=13, color=DARK_TEXT, spacing=Pt(3))

    # ══════════════════════════════════════════════════════════════
    # SLIDE 14 — Deployment Options
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _add_shape_bg(slide, WHITE)
    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
                 'Deployment Options', font_size=32, bold=True, color=DARK_TEXT)

    _add_card(slide, Inches(0.6), Inches(1.4), Inches(3.8), Inches(2.5),
              '💻  Power BI Desktop',
              'Double-click the .pbip file to open.\n\n'
              'Validate visuals, relationships, and measures in the native PBI experience.',
              PBI_YELLOW)

    _add_card(slide, Inches(4.8), Inches(1.4), Inches(3.8), Inches(2.5),
              '☁️  Power BI Service',
              'python migrate.py workbook.qvf \\\n'
              '  --deploy WORKSPACE_ID \\\n'
              '  --deploy-refresh\n\n'
              'Supports .pbix packaging + REST API upload.',
              AZURE_BLUE)

    _add_card(slide, Inches(9.0), Inches(1.4), Inches(3.8), Inches(2.5),
              '🏭  Microsoft Fabric',
              'Deploy Lakehouse, Dataflow, Notebook,\n'
              'SemanticModel, and Pipeline artifacts.\n\n'
              'Bundle deployment with endorsement.',
              RGBColor(0x00, 0x78, 0xD4))

    _add_card(slide, Inches(0.6), Inches(4.3), Inches(3.8), Inches(2.5),
              '🐳  Docker / REST API',
              'docker run -p 8000:8000 Qlik-to-pbi\n\n'
              'Headless migration via HTTP endpoints.\n'
              'Ideal for CI/CD pipelines.',
              ACCENT_PURPLE)

    _add_card(slide, Inches(4.8), Inches(4.3), Inches(3.8), Inches(2.5),
              '🏢  Multi-Tenant',
              'python migrate.py \\\n'
              '  --shared-model wb1.qvf wb2.qvf \\\n'
              '  --multi-tenant tenants.json\n\n'
              'Per-tenant connection string overrides.',
              GREEN)

    _add_card(slide, Inches(9.0), Inches(4.3), Inches(3.8), Inches(2.5),
              '🔄  Rolling Deploy',
              'Blue/green deployment with canary\nvalidation and auto-rollback.\n\n'
              'deploy_rolling() in pbi_deployer.py',
              QLIK_GREEN)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 15 — Key Stats
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _add_shape_bg(slide, DARK_BG)
    _add_textbox(slide, Inches(0.8), Inches(0.6), Inches(8), Inches(0.7),
                 'By the Numbers', font_size=36, bold=True, color=WHITE)

    stats = [
        ('175+', 'DAX\nConversions', QLIK_GREEN),
        ('75+', 'Visual Type\nMappings', PBI_YELLOW),
        ('40+', 'M Query\nTransforms', AZURE_BLUE),
        ('25', 'Data Source\nConnectors', GREEN),
        ('11', 'Intermediate\nJSON Files', ACCENT_PURPLE),
        ('2,000', 'Automated\nTests', RGBColor(0xEF, 0x44, 0x44)),
    ]
    for i, (num, label, color) in enumerate(stats):
        x = Inches(0.7) + Inches(i * 2.1)
        _add_stat_card(slide, x, Inches(1.8), Inches(1.9), Inches(1.3), num, label, color)

    _add_textbox(slide, Inches(0.8), Inches(3.8), Inches(11), Inches(0.5),
                 'Zero external dependencies · Python standard library only · Works on Python 3.12+',
                 font_size=16, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # Feature highlights
    features = [
        'Auto-generated Calendar table with Date Hierarchy',
        'Row-Level Security (RLS) from Section Access',
        'Set Analysis → CALCULATE with filter modifiers',
        'Cross-table RELATED() / LOOKUPVALUE() inference',
        'Aggr() decomposition → SUMX/COUNTX/AVERAGEX iterators',
        'Inter-record Above/Below/Peek → OFFSET/WINDOW',
    ]
    for i, feat in enumerate(features):
        col = i % 2
        row = i // 2
        x = Inches(1.0) + Inches(col * 5.5)
        y = Inches(4.5) + Inches(row * 0.45)
        _add_textbox(slide, x, y, Inches(5.5), Inches(0.4),
                     '✓  ' + feat, font_size=14, color=WHITE)

    # Migration results screenshot
    if os.path.isfile(IMG_MIGRATION_RESULTS):
        slide.shapes.add_picture(IMG_MIGRATION_RESULTS,
                                 Inches(0.7), Inches(6.0), Inches(11.8))

    # ══════════════════════════════════════════════════════════════
    # SLIDE 16 — Getting Started
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _add_shape_bg(slide, LIGHT_GRAY)
    _section_header(slide, 'Getting Started', '3 steps to your first migration')

    _add_card(slide, Inches(0.6), Inches(3.2), Inches(3.8), Inches(3.5),
              '1️⃣  Install',
              'git clone <repo-url>\ncd QlikToPowerBI\npython -m venv .venv\n'
              '.venv\\Scripts\\activate\n\n'
              'No pip install needed — zero dependencies.',
              AZURE_BLUE)

    _add_card(slide, Inches(4.8), Inches(3.2), Inches(3.8), Inches(3.5),
              '2️⃣  Migrate',
              'python migrate.py workbook.qvf\n\n'
              'Or assess first:\n'
              'python migrate.py workbook.qvf --assess\n\n'
              'Output lands in:\n'
              'artifacts/powerbi_projects/',
              QLIK_GREEN)

    _add_card(slide, Inches(9.0), Inches(3.2), Inches(3.8), Inches(3.5),
              '3️⃣  Validate & Deploy',
              'Open the .pbip in Power BI Desktop.\n\n'
              'Check:\n'
              '• Relationships in Model view\n'
              '• Visual accuracy vs Qlik\n'
              '• Measure values\n\n'
              'Then deploy to PBI Service or Fabric.',
              GREEN)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 17 — Multi-Agent Preceptorship Model
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _add_shape_bg(slide, WHITE)
    _add_textbox(slide, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
                 'Multi-Agent Preceptorship Model', font_size=32, bold=True, color=DARK_TEXT)
    _add_textbox(slide, Inches(0.8), Inches(1.0), Inches(10), Inches(0.4),
                 '10 AI agents collaborate following a Plan → Assign → Implement → Review cycle',
                 font_size=16, color=MID_GRAY)

    # Senior roles
    _add_card(slide, Inches(2.5), Inches(1.6), Inches(3.8), Inches(1.2),
              '🏗️  Tech Lead (Orchestrator)',
              'Architecture, planning, pipeline coordination, task decomposition',
              AZURE_BLUE)
    _add_card(slide, Inches(7.0), Inches(1.6), Inches(3.8), Inches(1.2),
              '🔍  Preceptor (Reviewer)',
              'Quality review, standards enforcement, pitfall detection, consistency',
              ACCENT_PURPLE)

    # Specialist agents
    agents = [
        ('Extractor', 'QVF/JSON parsing\n→ 11 JSON files', QLIK_GREEN),
        ('Converter', '175+ DAX conversions\n25 M connectors', QLIK_GREEN),
        ('Generator', 'TMDL + PBIR\n75+ visuals', AZURE_BLUE),
        ('Assessor', 'Readiness scoring\nStrategy advising', AZURE_BLUE),
        ('Merger', 'Multi-app merge\nThin reports', ACCENT_PURPLE),
        ('Deployer', 'Fabric / PBI Service\nAzure AD auth', GREEN),
        ('Tester', '2,000 tests\nRegression suites', RGBColor(0xEF, 0x44, 0x44)),
    ]
    for i, (name, desc, color) in enumerate(agents):
        x = Inches(0.4) + Inches(i * 1.8)
        _add_stat_card(slide, x, Inches(3.2), Inches(1.65), Inches(1.3),
                       name, '', color)
        _add_textbox(slide, x, Inches(4.6), Inches(1.65), Inches(1.0),
                     desc, font_size=11, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

    # Cycle description
    _add_textbox(slide, Inches(0.8), Inches(5.8), Inches(11.5), Inches(1.2),
                 'Each specialist follows the 4-step cycle:\n'
                 '1. Plan — outline changes, flag risks   '
                 '2. Assign — confirm ownership, declare handoffs   '
                 '3. Implement — smallest change that solves the problem   '
                 '4. Review — self-check: tests pass, contracts preserved, pitfalls avoided',
                 font_size=14, color=DARK_TEXT)

    # ══════════════════════════════════════════════════════════════
    # SLIDE 18 — Thank You
    # ══════════════════════════════════════════════════════════════
    slide = prs.slides.add_slide(blank)
    _add_shape_bg(slide, DARK_BG)
    _add_accent_bar(slide, QLIK_GREEN, Inches(0), Inches(3.2), Inches(13.333), Inches(0.06))
    _add_accent_bar(slide, PBI_YELLOW, Inches(0), Inches(3.28), Inches(13.333), Inches(0.06))
    _add_textbox(slide, Inches(1), Inches(2.0), Inches(11), Inches(0.8),
                 'Thank You', font_size=48, bold=True, color=WHITE,
                 alignment=PP_ALIGN.CENTER)
    _add_textbox(slide, Inches(1), Inches(4.0), Inches(11), Inches(0.5),
                 'python migrate.py your_app.qvf', font_size=20,
                 color=PBI_YELLOW, alignment=PP_ALIGN.CENTER, font_name='Cascadia Code')
    _add_textbox(slide, Inches(1), Inches(5.5), Inches(11), Inches(0.4),
                 'v9.1.0 · 2,000 tests · 175+ DAX conversions · 75+ visual mappings',
                 font_size=14, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

    # ── Save ──
    out = 'docs/qlik_to_powerbi_Migration_Guide.pptx'
    prs.save(out)
    print(f'✓ Presentation saved to {out}')
    print(f'  {len(prs.slides)} slides generated')


if __name__ == '__main__':
    build_presentation()
